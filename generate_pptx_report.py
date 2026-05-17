import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pptx(filename="BaoCao_ThuyetMinh_HeThong_LichCongTac_SIRMS.pptx"):
    prs = Presentation()
    # Kích thước màn hình Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Định nghĩa màu sắc chủ đạo
    c_primary = RGBColor(0x00, 0x61, 0xA3)   # Blue
    c_secondary = RGBColor(0x00, 0x83, 0x8F) # Teal
    c_dark = RGBColor(0x22, 0x22, 0x22)
    c_light = RGBColor(0xF4, 0xF5, 0xF7)
    c_white = RGBColor(0xFF, 0xFF, 0xFF)
    c_gray = RGBColor(0x66, 0x66, 0x66)
    
    # Hàm thêm slide bìa
    def add_cover_slide(title_text, subtitle_text, meta_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        
        # Thêm hình khối trang trí nền
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = c_light
        bg.line.color.rgb = c_light
        
        # Banner bên trái
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.5), Inches(7.5))
        banner.fill.solid()
        banner.fill.fore_color.rgb = c_primary
        banner.line.color.rgb = c_primary
        
        # Hộp văn bản tiêu đề
        txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10.5), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = subtitle_text
        p0.font.name = 'Arial'
        p0.font.size = Pt(20)
        p0.font.bold = True
        p0.font.color.rgb = c_secondary
        p0.space_after = Pt(20)
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.name = 'Arial'
        p1.font.size = Pt(32)
        p1.font.bold = True
        p1.font.color.rgb = c_primary
        p1.space_after = Pt(40)
        
        # Hộp văn bản metadata
        txBox_meta = slide.shapes.add_textbox(Inches(1.5), Inches(5.0), Inches(10.5), Inches(2.0))
        tf_meta = txBox_meta.text_frame
        tf_meta.word_wrap = True
        
        for line in meta_text:
            p = tf_meta.add_paragraph()
            p.text = line
            p.font.name = 'Arial'
            p.font.size = Pt(16)
            p.font.color.rgb = c_dark
            p.space_after = Pt(6)
            
    # Hàm thêm slide nội dung chuẩn
    def add_content_slide(title_text, bullets):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = c_primary
        header.line.color.rgb = c_primary
        
        # Tiêu đề slide
        txBox_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
        tf_t = txBox_t.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = c_white
        
        # Nội dung
        txBox_c = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
        tf_c = txBox_c.text_frame
        tf_c.word_wrap = True
        
        for i, b in enumerate(bullets):
            if i == 0:
                p = tf_c.paragraphs[0]
            else:
                p = tf_c.add_paragraph()
            p.text = b
            p.font.name = 'Arial'
            p.font.size = Pt(18)
            p.font.color.rgb = c_dark
            p.space_after = Pt(16)
            
    # Hàm thêm slide bảng
    def add_table_slide(title_text, table_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # Header banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.2))
        header.fill.solid()
        header.fill.fore_color.rgb = c_primary
        header.line.color.rgb = c_primary
        
        txBox_t = slide.shapes.add_textbox(Inches(0.8), Inches(0.2), Inches(11.5), Inches(0.8))
        tf_t = txBox_t.text_frame
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.name = 'Arial'
        p_t.font.size = Pt(26)
        p_t.font.bold = True
        p_t.font.color.rgb = c_white
        
        # Bảng
        rows = len(table_data)
        cols = len(table_data[0])
        left = Inches(0.8)
        top = Inches(1.6)
        width = Inches(11.7)
        height = Inches(5.2)
        
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        col_widths = [Inches(2.0), Inches(1.5), Inches(2.2), Inches(3.2), Inches(2.8)]
        for j, w in enumerate(col_widths):
            table.columns[j].width = w
            
        for i, row in enumerate(table.rows):
            for j, cell in enumerate(row.cells):
                cell.text = table_data[i][j]
                p = cell.text_frame.paragraphs[0]
                p.font.name = 'Arial'
                p.font.size = Pt(13)
                if i == 0:
                    p.font.bold = True
                    p.font.color.rgb = c_white
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = c_primary
                else:
                    p.font.color.rgb = c_dark
                    cell.fill.solid()
                    if i % 2 == 1:
                        cell.fill.fore_color.rgb = c_white
                    else:
                        cell.fill.fore_color.rgb = c_light

    # 1. Slide bìa
    add_cover_slide(
        "BÁO CÁO THUYẾT MINH HỆ THỐNG LỊCH CÔNG TÁC KẾT HỢP THEO DÕI THIẾT BỊ DỰA TRÊN TẬN DỤNG HỆ THỐNG SIRMS",
        "BÁO CÁO THUYẾT MINH & HƯỚNG DẪN TRIỂN KHAI",
        [
            "Đơn vị triển khai: Đài DVOR/DME Tuy Hòa",
            "Tác giả / Hỗ trợ kỹ thuật: Nguyễn Hoàng Hải",
            "Phiên bản tài liệu: 2.5 (Tích hợp PWA & Sinh trắc học WebAuthn)",
            "Thời gian cập nhật: Tháng 5 / 2026"
        ]
    )

    # 2. Slide Tổng quan
    add_content_slide(
        "1. TỔNG QUAN KIẾN TRÚC & KHẢ NĂNG MỞ RỘNG",
        [
            "• Giải pháp số hóa toàn diện: Kết hợp giữa Quản lý lịch trực điều hành nội bộ và Giám sát thiết bị từ xa (SIRMS).",
            "• Kiến trúc Serverless / Micro-services: Tách bạch hoàn toàn giữa Data Collector (trạm nội bộ) và Cloud Backend.",
            "• Khả năng linh hoạt vượt trội: Dễ dàng áp dụng cho các đơn vị đài trạm CNS khác mà không cần sửa đổi lõi hệ thống.",
            "• Tối ưu hóa chi phí (Zero Server Cost): Hoạt động 100% trên hạ tầng Cloudflare (Pages, Workers, D1, KV), không tốn chi phí mua sắm hay duy trì server."
        ]
    )

    # 3. Slide Chức năng 1
    add_content_slide(
        "2.1. QUẢN LÝ LỊCH CÔNG TÁC & CÔNG VIỆC HÀNG NGÀY",
        [
            "• Quản lý Lịch trực ban theo ca: Phân chia ca Hành chính (X), Ca 1 (X1), Ca 2 (X2), Ca 3 (Đ) với màu sắc trực quan. Hỗ trợ import tự động từ file Excel.",
            "• Quản lý Công việc hàng ngày (Tasks): Thêm, sửa, xóa công việc theo khung giờ, gán nhân sự tham gia và theo dõi tiến độ.",
            "• Công việc định kỳ tự động (Recurring Tasks): Hỗ trợ thiết lập việc lặp lại theo ngày cố định, dải ngày, hoặc Thứ 6 ưu tiên (Friday Priority).",
            "• Ghi chú nhanh (Daily Notes): Công cụ tiện lợi giúp ca trực ghi nhận sự kiện phát sinh và bàn giao ca dễ dàng."
        ]
    )

    # 4. Slide Chức năng 2
    add_content_slide(
        "2.2. GIÁM SÁT THÔNG SỐ ĐÀI DẪN ĐƯỜNG (SIRMS)",
        [
            "• Bảng thông số trực quan (Navaid Widget): Hiển thị thời gian thực các tham số VOR (Azimuth, Mod 30Hz, Mod 9960Hz, Deviation, RF Level) và DME (Delay, Spacing, Tx Power, ERP, Efficiency).",
            "• Biểu đồ phân tích lịch sử: Tích hợp ApexCharts theo dõi biến động thông số trong 48 giờ hoặc 7 ngày qua, phát hiện sớm suy giảm chất lượng.",
            "• Hệ thống Cảnh báo Ngưỡng Giới Hạn: Thiết lập ngưỡng Max/Min cho từng tham số. Phân loại tự động: Cảnh báo (Alert), Báo động (Alarm), Báo động Khẩn cấp (Critical Alarm)."
        ]
    )

    # 5. Slide Chức năng 3
    add_content_slide(
        "2.3. BẢO MẬT ĐA TẦNG & PWA PUSH NOTIFICATIONS",
        [
            "• Phân quyền người dùng chặt chẽ: Chia làm 2 cấp quyền ADMIN (Quản trị) và VIEWER (Người xem).",
            "• Xác thực Fallback & Chống Brute-force: Đăng nhập bằng mã PIN băm (hash) kết hợp tự động kích hoạt Captcha toán học khi nhập sai nhiều lần.",
            "• Đăng nhập Sinh trắc học (WebAuthn / Passkeys): Đăng nhập siêu nhanh bằng Face ID / Touch ID / Vân tay trên PWA. Xác thực chữ ký ECDSA P-256 qua Web Crypto API.",
            "• Thông báo đẩy PWA (Web Push Notifications): Gửi cảnh báo thời gian thực tới điện thoại của kỹ thuật viên khi thông số vượt ngưỡng."
        ]
    )

    # 6. Slide Quy trình hoạt động Python
    add_content_slide(
        "3. QUY TRÌNH THU THẬP DỮ LIỆU LOCAL (PYTHON SCRIPT)",
        [
            "• Service chạy nền 24/7 (Daemon): Hoạt động độc lập trên PC nội bộ của đài trạm.",
            "• Vòng lặp Polling đa luồng: Tự động truy vấn thiết bị phần cứng (192.168.105.25/alldata.json hoặc SNMP get) mỗi 5 giây.",
            "• Chuẩn hóa dữ liệu tại nguồn: Hàm safe_float() bóc tách số và chia tỷ lệ (10/100) để đưa về đúng đơn vị đo lường thực tế (Ví dụ: 249.46°, 29.4 Hz, 100%).",
            "• Phục vụ qua HTTP Flask: Mở port 5050 cung cấp các endpoint /navaid và /health cho Cloudflare Tunnel kết nối."
        ]
    )

    # 7. Slide Bảng tham số
    add_table_slide(
        "BẢNG CÁC THAM SỐ CHUẨN HÓA TRONG PYTHON",
        [
            ["Tên Tham Số", "Khối Thiết Bị", "Dữ Liệu Thô (Raw)", "Hàm Chuẩn Hóa", "Kết Quả Chuẩn Hóa"],
            ["Azimuth", "DVOR", "24946 degree", "safe_float(..., 100)", "249.46 °"],
            ["Mod30Hz", "DVOR", "294 Hz", "safe_float(..., 10)", "29.4 Hz"],
            ["Mod9960Hz", "DVOR", "307 Hz", "safe_float(..., 10)", "30.7 Hz"],
            ["Deviation", "DVOR", "1615", "safe_float(..., 100)", "16.15"],
            ["RFLevel", "DVOR", "-9 dBm", "safe_float(..., 10)", "-0.9 dBm"],
            ["Delay", "DME", "5001 us", "safe_float(..., 100)", "50.01 µs"],
            ["Spacing", "DME", "1198 us", "safe_float(..., 100)", "11.98 µs"],
            ["TxPower", "DME", "1061 Watts", "safe_float(..., 1)", "1061 W"],
            ["ERP", "DME", "15", "safe_float(..., 10)", "1.5 dB"],
            ["Efficiency", "DME", "1000%", "safe_float(..., 10)", "100 %"],
        ]
    )

    # 8. Slide Cấu hình Cloudflare
    add_content_slide(
        "4. HƯỚNG DẪN CẤU HÌNH HỆ SINH THÁI CLOUDFLARE",
        [
            "• Cloudflare Tunnel (cloudflared): Tạo đường hầm bảo mật trỏ localhost:5050 ra sirms-api.hainh.io.vn. Tích hợp WAF rule kiểm tra Header authorization để chặn truy cập trái phép.",
            "• Cloudflare D1 Database (SQL Serverless): Khởi tạo các bảng tasks, shift_schedules, sirms_data, navaid_limits, push_subscriptions, webauthn_credentials.",
            "• Cloudflare KV Namespace: Bộ nhớ đệm tốc độ cao lưu trữ Rate Limiting, Captcha và WebAuthn Challenges.",
            "• sirms-worker & Cron Trigger: Cấu hình chạy tự động mỗi 15 phút (0/15 * * * *), lấy dữ liệu từ Tunnel, lưu D1, kiểm tra ngưỡng và bắn Push Notification.",
            "• Cloudflare Pages & K1 Nameserver: Triển khai PWA React Frontend, quản lý DNS tập trung."
        ]
    )

    # 9. Slide Kết luận
    add_content_slide(
        "5. TỔNG KẾT & KHUYẾN NGHỊ NHÂN RỘNG",
        [
            "• Hiệu quả Kỹ thuật: Hệ thống hoạt động ổn định 24/7, bảo mật tuyệt đối nhờ kiến trúc Serverless và xác thực đa tầng. Trải nghiệm PWA hiện đại như ứng dụng Native.",
            "• Hiệu quả Kinh tế: Tận dụng tối đa hạ tầng miễn phí của Cloudflare, tiết kiệm 100% chi phí mua sắm và vận hành máy chủ.",
            "• Khuyến nghị Nhân rộng: Với khả năng linh hoạt cao trong việc thay đổi nguồn thu thập dữ liệu, chỉ cần thay đổi giao diện và cấu hình IP của thiết bị SNMP là có thể áp dụng cho các đài trạm CNS khác trong Trung tâm Bảo đảm Kỹ thuật."
        ]
    )

    prs.save(filename)
    print(f"Successfully generated PPTX: {filename}")

if __name__ == "__main__":
    create_pptx()
